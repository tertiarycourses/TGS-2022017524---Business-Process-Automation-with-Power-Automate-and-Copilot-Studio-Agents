#!/usr/bin/env python3
"""One-time WSQ revision of courseware/Business Process Automation with Power Automate and Copilot Studio Agents-v3.pptx.

- Rebrands the cover slide to the new course title + org + UEN + TGS code.
- Inserts the WSQ admin slides (Digital Attendance, About the Trainer x2,
  Learner Introduction, Ground Rules, Lesson Plan, Briefing for Assessment,
  Assessment & Funding, Assessment Flow, Access the Hands-On Labs) after the cover.
- Inserts the end-of-deck assessment block (LMS, Assessment, Assessment Flow,
  Digital Attendance) before the Thank You slide.
- Adds the copyright footer to every slide.

Final order (108 slides):
  1 cover | 2-6 admin A-E | 7-9 old 2-4 | 10-14 admin F-J | 15-103 old 5-93
  | 104-107 end block K-N | 108 old 94 (Thank You)
"""
import copy, math, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# script lives at .claude/skills/course-slides/ — repo root is 3 levels up
REPO = os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", ".."))
DECK = os.path.join(REPO, "courseware", "Business Process Automation with Power Automate and Copilot Studio Agents-v3.pptx")

TITLE = "Business Process Automation with Power Automate and Copilot Studio Agents"
SHORT = "Business Process Automation with Power Automate & Copilot Studio Agents"
CODE = "TGS-2022017524"
REPO_URL = ("https://github.com/tertiarycourses/"
            "TGS-2022017524---Business-Process-Automation-with-Power-Automate-and-Copilot-Studio-Agents")

BLUE = RGBColor(0x1F, 0x6F, 0xEB); TEAL = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B); VIOLET = RGBColor(0x7C, 0x3A, 0xED)
INK = RGBColor(0x16, 0x1B, 0x26); GREY = RGBColor(0x5B, 0x63, 0x72)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
PALETTE = [BLUE, TEAL, VIOLET, AMBER]

prs = Presentation(DECK)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------- helpers
def slide(): return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, color):
    sp = s.shapes.add_shape(1, x, y, w, h); sp.fill.solid()
    sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def oval(s, x, y, w, h, color):
    sp = s.shapes.add_shape(9, x, y, w, h); sp.fill.solid()
    sp.fill.fore_color.rgb = color; sp.line.fill.background(); sp.shadow.inherit = False
    return sp

def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, ln in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        for t, sz, col, bold in ln:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.bold = bold
            r.font.color.rgb = col; r.font.name = "Arial"
    return tb

def bullets(s, x, y, w, h, items, size=18, gap=10):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(gap)
        lvl = it[1] if isinstance(it, tuple) else 0
        text = it[0] if isinstance(it, tuple) else it
        r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + text
        r.font.size = Pt(size if lvl == 0 else size - 2)
        r.font.color.rgb = INK if lvl == 0 else GREY; r.font.name = "Arial"
    return tb

def head(s, title, kicker=None, kcolor=BLUE):
    rect(s, 0, 0, SW, SH, WHITE); rect(s, 0, 0, Inches(0.28), Inches(1.55), kcolor)
    if kicker:
        txt(s, Inches(0.85), Inches(0.5), Inches(11.6), Inches(0.4), [[(kicker, 14, kcolor, True)]])
    txt(s, Inches(0.85), Inches(0.9), Inches(11.9), Inches(0.9), [[(title, 29, INK, True)]])
    rect(s, Inches(0.85), Inches(1.7), Inches(11.63), Inches(0.02), LINE)
    return s

def content(title, items, kicker=None, size=18):
    s = head(slide(), title, kicker)
    bullets(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(4.9), items, size=size)
    return s

def two_col(title, left, right, kicker=None, lhead="", rhead=""):
    s = head(slide(), title, kicker)
    rect(s, Inches(0.85), Inches(1.95), Inches(5.7), Inches(4.7), LIGHT)
    rect(s, Inches(6.95), Inches(1.95), Inches(5.55), Inches(4.7), LIGHT)
    if lhead: txt(s, Inches(1.1), Inches(2.15), Inches(5.2), Inches(0.4), [[(lhead, 16, BLUE, True)]])
    if rhead: txt(s, Inches(7.2), Inches(2.15), Inches(5.0), Inches(0.4), [[(rhead, 16, TEAL, True)]])
    bullets(s, Inches(1.1), Inches(2.7), Inches(5.2), Inches(3.8), left, size=15, gap=8)
    bullets(s, Inches(7.2), Inches(2.7), Inches(5.05), Inches(3.8), right, size=15, gap=8)
    return s

def tile_grid(title, items, kicker=None, cols=2, size=15, icons=None, accent=BLUE):
    s = head(slide(), title, kicker, kcolor=accent)
    n = len(items); rows = math.ceil(n / cols)
    X0 = Inches(0.85); Y0 = Inches(1.95); TOTW = Inches(11.63); AREAH = Inches(4.78)
    gx = Inches(0.3); gy = Inches(0.26)
    cw = int((TOTW - gx * (cols - 1)) / cols); ch = int((AREAH - gy * (rows - 1)) / rows)
    bd = Inches(0.6)
    for i, it in enumerate(items):
        r = i // cols; c = i % cols
        x = int(X0 + (cw + gx) * c); y = int(Y0 + (ch + gy) * r); col = PALETTE[i % len(PALETTE)]
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, Inches(0.1), ch, col)
        oval(s, x + Inches(0.28), int(y + ch / 2 - bd / 2), bd, bd, col)
        ic = icons[i] if icons else str(i + 1)
        txt(s, x + Inches(0.28), int(y + ch / 2 - bd / 2), bd, bd, [[(ic, 19, WHITE, True)]],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx = x + Inches(1.08); tw = cw - Inches(1.32)
        if isinstance(it, tuple):
            txt(s, tx, int(y + Inches(0.14)), tw, int(ch - Inches(0.2)),
                [[(it[0], size + 2, INK, True)], [(it[1], size - 2, GREY, False)]],
                anchor=MSO_ANCHOR.MIDDLE, space=3)
        else:
            txt(s, tx, int(y + Inches(0.1)), tw, int(ch - Inches(0.16)),
                [[(it, size, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
    return s

def flow_h(title, steps, kicker=None, color=BLUE):
    s = head(slide(), title, kicker, kcolor=color)
    n = len(steps); X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.34)
    cw = int((TOTW - gap * (n - 1)) / n); y = Inches(2.55); ch = Inches(3.15); bd = Inches(0.82)
    for i, st in enumerate(steps):
        x = int(X0 + (cw + gap) * i)
        rect(s, x, y, cw, ch, LIGHT); rect(s, x, y, cw, Inches(0.1), color)
        oval(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd, color)
        txt(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd,
            [[(str(i + 1), 30, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        txt(s, x + Inches(0.16), int(y + Inches(1.55)), cw - Inches(0.32), int(ch - Inches(1.7)),
            [[(st, 14, INK, False)]], align=PP_ALIGN.CENTER)
        if i < n - 1:
            txt(s, int(x + cw - Inches(0.04)), int(y + ch / 2 - Inches(0.3)),
                int(gap + Inches(0.08)), Inches(0.6),
                [[("▶", 15, color, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s

def trainer_slide(kicker, name, role, rows, initials, accent=BLUE):
    s = head(slide(), "About the Trainer", kicker, kcolor=accent)
    lx = Inches(0.85); lw = Inches(3.65)
    rect(s, lx, Inches(1.95), lw, Inches(4.7), LIGHT); rect(s, lx, Inches(1.95), lw, Inches(0.12), accent)
    bd = Inches(1.7); ax = int(lx + (lw - bd) / 2)
    oval(s, ax, Inches(2.5), bd, bd, accent)
    txt(s, ax, Inches(2.5), bd, bd, [[(initials, 44, WHITE, True)]],
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, lx + Inches(0.15), Inches(4.55), lw - Inches(0.3), Inches(0.6),
        [[(name, 21, INK, True)]], align=PP_ALIGN.CENTER)
    txt(s, lx + Inches(0.15), Inches(5.2), lw - Inches(0.3), Inches(1.2),
        [[(role, 13, GREY, False)]], align=PP_ALIGN.CENTER)
    rx = Inches(4.9); rw = Inches(7.6); ry = Inches(1.95); rh = Inches(4.7)
    n = len(rows); gy = Inches(0.2); th = int((rh - gy * (n - 1)) / n)
    for i, (label, val) in enumerate(rows):
        y = int(ry + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
        rect(s, rx, y, rw, th, LIGHT); rect(s, rx, y, Inches(0.1), th, col)
        vruns = [(val, 14, INK, False)] if val else [("____________________________________________", 13, LINE, False)]
        txt(s, rx + Inches(0.32), y, rw - Inches(0.6), th,
            [[(label.upper(), 11, col, True)], vruns], anchor=MSO_ANCHOR.MIDDLE, space=3)
    return s

# ---------------------------------------------------------------- 1. cover
cov = prs.slides[0]
tsh, ssh = cov.shapes[0], cov.shapes[1]
tsh.text_frame.paragraphs[0].runs[0].text = "Business Process Automation"
# subtitle placeholder: rebuild its two paragraphs' first runs
sp = ssh.text_frame.paragraphs
sp[0].runs[0].text = "with Power Automate and Copilot Studio Agents"
for para in sp[1:]:
    if para.runs:
        para.runs[0].text = f"WSQ Course Code: {CODE}  ·  3-Day Hands-On Training  ·  Modules 1–5"
        break
txt(cov, Inches(1.67), Inches(6.15), Inches(10.0), Inches(0.85),
    [[("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W", 14, GREY, False)],
     [("Version 2.0  ·  www.tertiarycourses.com.sg  ·  https://lms-tms.tertiaryinfotech.com/", 11, GREY, False)]],
    align=PP_ALIGN.CENTER, space=4)

# ---------------------------------------------------------------- 2. new slides
new_slides = []
def reg(s): new_slides.append(s); return s

ATTEND = ["It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
          "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
          "Scan the QR code with your mobile phone camera and submit your attendance.",
          "A minimum of 75% attendance is required to be eligible for assessment and funding."]
FLOW = ["TRAQOM survey — scan the QR code on the LMS",
        "Assessment digital attendance — scan the SSG QR",
        "Sit WA (SAQ) then PP — open book",
        "Submit your answers on the LMS",
        "Sign the Assessment Summary Record"]

# A
reg(content("Digital Attendance (Mandatory)", ATTEND, kicker="TRAQOM · SSG DIGITAL ATTENDANCE"))
# B
reg(trainer_slide("YOUR TRAINER · GENERAL", "Your Trainer",
    "General Trainer template —\nto be completed by the trainer",
    [("Name", ""), ("Title / Designation", ""), ("Qualifications", ""),
     ("Areas of expertise", ""), ("Training & industry experience", ""), ("Contact", "")],
    initials="?", accent=GREY))
# C
reg(trainer_slide("YOUR TRAINER", "Dr. Alfred Ang",
    "Principal Trainer\nTertiary Infotech Academy Pte. Ltd.",
    [("Role", "Principal Trainer, Tertiary Infotech Academy Pte. Ltd."),
     ("Certification", "Microsoft Power Platform, AI and data analytics — 20+ years of training experience."),
     ("Delivers", "WSQ courses on AI, business automation, data science and software engineering."),
     ("Founder", "Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
    initials="AA", accent=BLUE))
# D
reg(content("Let's Know Each Other", [
    "Your name and organisation / role.",
    "Your experience with Power Automate, Copilot Studio or other automation tools (if any).",
    "One repetitive business process you would love to automate after this course."],
    kicker="ICE-BREAKER"))
# E
reg(tile_grid("Ground Rules", [
    "Set your mobile phone to silent mode.", "Participate actively — no question is too small.",
    "Mutual respect: agree to disagree.", "One conversation at a time.",
    "Be punctual; return from breaks on time.", "75% attendance is required."],
    kicker="HOUSEKEEPING", cols=2, size=15))
# F
reg(two_col("Lesson Plan — 3 Days, 9:00am–6:00pm", [
    ("Day 1 — Foundations & Power Automate", 0),
    ("Modules 1–2: workflow concepts, Power Automate", 1),
    ("Labs 0–5: email, Excel logging, approval, scheduled, form flows", 1),
    ("Day 2 — Business Agents with Copilot Studio", 0),
    ("Module 3: agents, prompts, agent + flow", 1),
    ("Labs 6–11: first agent, knowledge (RAG), tools, agent flows", 1)],
    [("Day 3 — End-to-End Automation & Workshop", 0),
     ("Module 4: orchestration · Module 5: workshop", 1),
     ("Labs 12–15: end-to-end business workflows", 1),
     ("Lab 16: capstone workshop — build & present", 1),
     ("Assessment (Day 3, 4:00–6:00pm)", 0),
     ("Written Assessment 1 hr + Practical Performance 1 hr", 1),
     ("Daily timing", 0),
     ("9:00am–6:00pm · 1-hour lunch · tea breaks within", 1)],
    kicker="SCHEDULE", lhead="Days 1–2", rhead="Day 3 & assessment"))
# G
reg(content("Briefing for Assessment", [
    "Place phones and other materials under the table or on the floor.",
    "No photos or recording of assessment scripts.",
    "No discussion during the assessment.",
    "Use a black/blue pen for hard-copy assessments.",
    "No liquid paper / correction tape.",
    "Scripts are collected when time is up."], kicker="BEFORE THE ASSESSMENT"))
# H
reg(content("Assessment & Funding", [
    "Written Assessment (SAQ) — 1 hour, open book. Tests the knowledge from Modules 1–4.",
    "Practical Performance (PP) — 1 hour, open book. You build a working flow + agent, as in the labs.",
    "Both are taken on Day 3, 4:00–6:00pm, and submitted on the LMS.",
    "WSQ funding requires 75% attendance AND passing the assessment (Competent).",
    "Courseware and the assessment are on the LMS: https://lms-tms.tertiaryinfotech.com/"],
    kicker="FINAL ASSESSMENT"))
# I
reg(flow_h("Assessment Flow", FLOW, kicker="ON ASSESSMENT DAY"))
# J — Access the Hands-On Labs (clickable repo link + two option cards)
s = head(slide(), "Access the Hands-On Labs", kicker="COURSE REPOSITORY")
tb = txt(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(0.5),
         [[("All 17 labs + the Learner Guide live in the course GitHub repository:", 16, INK, False)]])
lb = s.shapes.add_textbox(Inches(0.85), Inches(2.45), Inches(11.6), Inches(0.5))
lp = lb.text_frame.paragraphs[0]; lr = lp.add_run()
lr.text = REPO_URL.replace("https://", "")
lr.font.size = Pt(15); lr.font.bold = True; lr.font.name = "Arial"; lr.font.color.rgb = BLUE
lr.hyperlink.address = REPO_URL
for x, col, htxt, lines in [
        (Inches(0.85), BLUE, "Option A — Clone with Git",
         ["Install Git (git-scm.com) if needed.", "git clone <repo URL>.git",
          "Open the labs/ folder in your editor or browser."]),
        (Inches(6.95), TEAL, "Option B — Download the ZIP",
         ["Open the repo URL in your browser.", "Code → Download ZIP, then unzip.",
          "Open labs/ and follow each lab's index.md."])]:
    rect(s, x, Inches(3.2), Inches(5.55), Inches(3.3), LIGHT)
    rect(s, x, Inches(3.2), Inches(5.55), Inches(0.12), col)
    txt(s, x + Inches(0.25), Inches(3.45), Inches(5.05), Inches(0.5), [[(htxt, 17, col, True)]])
    bullets(s, x + Inches(0.25), Inches(4.05), Inches(5.05), Inches(2.3), lines, size=14, gap=8)
reg(s)
# K — LMS (end block)
reg(content("Courseware & Assessment on the LMS", [
    "Download the slides and the Learner Guide from the LMS for the open-book assessment.",
    "Portal: https://lms-tms.tertiaryinfotech.com/",
    "Submit your Written Assessment and Practical Performance answers on the LMS.",
    "Complete the TRAQOM survey via the QR code on the LMS."], kicker="COURSE PORTAL"))
# L
reg(content("Assessment", [
    "Written Assessment (SAQ) — 1 hour. Practical Performance (PP) — 1 hour.",
    "Open book: slides, Learner Guide and approved materials only.",
    "Remember to take the Assessment digital attendance (TRAQOM).",
    "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],
    kicker="FINAL ASSESSMENT"))
# M
reg(flow_h("Assessment Flow", FLOW, kicker="ON ASSESSMENT DAY"))
# N
reg(content("Digital Attendance (Mandatory)", ATTEND, kicker="TRAQOM · SSG DIGITAL ATTENDANCE"))

# ---------------------------------------------------------------- 3. reorder
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)          # 94 original + 14 new (appended in order A..N)
old, new = ids[:94], ids[94:]
A, B, C, D, E, F, G, H, I, J, K, L, M, N = new
order = ([old[0], A, B, C, D, E] + old[1:4] + [F, G, H, I, J]
         + old[4:93] + [K, L, M, N] + [old[93]])
assert len(order) == 108
for el in ids: sldIdLst.remove(el)
for el in order: sldIdLst.append(el)

# ---------------------------------------------------------------- 4. footers
for i, s in enumerate(prs.slides, 1):
    txt(s, Inches(0.4), Inches(7.08), Inches(9.2), Inches(0.32),
        [[(f"{SHORT}  ·  {CODE}  ·  © 2026 Tertiary Infotech Academy Pte Ltd", 9, GREY, False)]])
    txt(s, Inches(12.35), Inches(7.08), Inches(0.62), Inches(0.32),
        [[(str(i), 9, GREY, False)]], align=PP_ALIGN.RIGHT)

prs.save(DECK)
print(f"Saved {DECK} with {len(prs.slides._sldIdLst)} slides")
