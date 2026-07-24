#!/usr/bin/env python3
"""Rebuild key concept and lab-roadmap slides as visual teaching diagrams.

The existing 86-slide course deck is edited in place. Slide count and lesson-plan
mapping remain stable while the progression is clarified:

instant -> scheduled -> automated -> approval -> HTTP -> agent -> RAG ->
Teams/website -> agent flow -> prompt flow.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DECK = ROOT / "courseware" / (
    "Business Process Automation with Power Automate and Copilot Studio Agents-v3.pptx"
)

BLUE = RGBColor(0x1F, 0x6F, 0xEB)
TEAL = RGBColor(0x10, 0x8A, 0x73)
GREEN = RGBColor(0x16, 0x84, 0x5B)
AMBER = RGBColor(0xC7, 0x76, 0x00)
VIOLET = RGBColor(0x6D, 0x3F, 0xD2)
RED = RGBColor(0xC2, 0x41, 0x3A)
INK = RGBColor(0x16, 0x1B, 0x26)
GREY = RGBColor(0x5B, 0x63, 0x72)
LINE = RGBColor(0xD7, 0xE0, 0xEA)
LIGHT = RGBColor(0xF5, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALETTE = [BLUE, TEAL, VIOLET, AMBER, GREEN]


prs = Presentation(DECK)
SW, SH = prs.slide_width, prs.slide_height


def clear(slide):
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        x,
        y,
        w,
        h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    return shape


def text(
    slide,
    x,
    y,
    w,
    h,
    value,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.MIDDLE,
    margin=0.08,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    r.font.name = "Arial"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return tb


def title(slide, heading, kicker):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), Inches(1.48))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    text(slide, Inches(0.72), Inches(0.38), Inches(11.8), Inches(0.28), kicker.upper(), 12, BLUE, True)
    text(slide, Inches(0.72), Inches(0.72), Inches(11.8), Inches(0.62), heading, 29, INK, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.5), Inches(11.85), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def footer(slide, slide_no):
    text(
        slide,
        Inches(0.72),
        Inches(7.1),
        Inches(11.85),
        Inches(0.2),
        f"Business Process Automation with Power Automate and Copilot Studio Agents   |   {slide_no}",
        8,
        GREY,
        False,
        PP_ALIGN.RIGHT,
    )


def card(slide, x, y, w, h, label, body, accent=BLUE, number=None, body_size=14, label_size=17):
    box(slide, x, y, w, h, LIGHT, LINE)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.09), h)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()
    if number is not None:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.22), y + Inches(0.22), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.fill.background()
        text(slide, x + Inches(0.22), y + Inches(0.22), Inches(0.5), Inches(0.5), str(number), 16, WHITE, True, PP_ALIGN.CENTER)
        lx = x + Inches(0.84)
        lw = w - Inches(1.05)
    else:
        lx = x + Inches(0.28)
        lw = w - Inches(0.52)
    text(slide, lx, y + Inches(0.16), lw, Inches(0.45), label, label_size, accent, True)
    text(slide, x + Inches(0.28), y + Inches(0.7), w - Inches(0.52), h - Inches(0.82), body, body_size, GREY)


def arrow(slide, x, y, w=0.42, h=0.38, color=BLUE, direction="right"):
    kind = MSO_SHAPE.CHEVRON if direction == "right" else MSO_SHAPE.DOWN_ARROW
    shape = slide.shapes.add_shape(kind, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def flow(slide, y, nodes, labels=None, colors=None, x0=0.72, total=11.85, height=1.15):
    labels = labels or [""] * (len(nodes) - 1)
    colors = colors or [BLUE] * len(nodes)
    gap = 0.46
    node_w = (total - gap * (len(nodes) - 1)) / len(nodes)
    for i, node in enumerate(nodes):
        x = Inches(x0 + i * (node_w + gap))
        box(slide, x, Inches(y), Inches(node_w), Inches(height), WHITE, colors[i])
        text(slide, x + Inches(0.08), Inches(y + 0.08), Inches(node_w - 0.16), Inches(height - 0.16), node, 13, INK, True, PP_ALIGN.CENTER)
        if i < len(nodes) - 1:
            ax = x + Inches(node_w + 0.06)
            arrow(slide, ax, Inches(y + height / 2 - 0.17), Inches(0.3), Inches(0.34), colors[i])
            if labels[i]:
                text(slide, x + Inches(node_w - 0.12), Inches(y - 0.28), Inches(gap + 0.24), Inches(0.25), labels[i], 9, GREY, False, PP_ALIGN.CENTER)


def rebuild(slide_no, heading, kicker, draw):
    slide = prs.slides[slide_no - 1]
    clear(slide)
    title(slide, heading, kicker)
    draw(slide)
    footer(slide, slide_no)


def draw_course_outcomes(s):
    flow(s, 2.05, ["Power Automate\nflows", "Copilot Studio\nagent", "Teams + website\nchannels", "Agent tools\nand prompts"], ["add AI", "publish", "take action"], [BLUE, VIOLET, TEAL, AMBER], height=1.35)
    card(s, Inches(0.72), Inches(4.05), Inches(3.65), Inches(1.8), "Day 1 — automate", "Build reliable trigger → action flows without an AI agent.", BLUE)
    card(s, Inches(4.65), Inches(4.05), Inches(3.65), Inches(1.8), "Day 2 — converse", "Create a grounded agent that understands and confirms user intent.", VIOLET)
    card(s, Inches(8.58), Inches(4.05), Inches(3.99), Inches(1.8), "Integrate — act", "Let the agent call deterministic or prompt-based Power Automate tools.", TEAL)


def draw_journey(s):
    items = [
        ("1–2", "Instant", "Run manually; email and Excel"),
        ("3", "Scheduled", "Recurrence and time zone"),
        ("4", "Automated", "Microsoft Forms event"),
        ("5", "Approval", "Human decision and branch"),
        ("6A/B", "HTTP", "External JSON request/response"),
        ("7A", "Agent", "Instructions and safe behaviour"),
        ("7B", "RAG", "Retrieve approved knowledge"),
        ("8", "Channels", "Publish to Teams and website"),
        ("9", "Agent flow", "Deterministic tool call"),
        ("10", "Prompt flow", "Guarded AI generation"),
    ]
    for i, (num, label, body) in enumerate(items):
        row, col = divmod(i, 5)
        x = Inches(0.72 + col * 2.42)
        y = Inches(1.9 + row * 2.35)
        card(s, x, y, Inches(2.12), Inches(1.92), label, body, PALETTE[i % 5], num, 11, 12)
    text(s, Inches(0.72), Inches(6.54), Inches(11.85), Inches(0.34), "Each lab reuses the previous capability; complexity increases one controlled step at a time.", 14, INK, True, PP_ALIGN.CENTER)


def draw_trigger_spectrum(s):
    types = [
        ("Instant", "A person presses Run", "Labs 1–2", BLUE),
        ("Scheduled", "A clock reaches a timetable", "Lab 3", TEAL),
        ("Automated", "An event occurs", "Lab 4", VIOLET),
        ("Agent flow", "A Copilot agent calls a tool", "Labs 9–10", AMBER),
    ]
    for i, (label, body, lab, color) in enumerate(types):
        x = Inches(0.72 + i * 3.0)
        card(s, x, Inches(1.9), Inches(2.72), Inches(2.1), label, f"{body}\n\n{lab}", color, i + 1, 13)
    card(s, Inches(1.55), Inches(4.55), Inches(4.9), Inches(1.55), "Human-in-the-loop pattern", "Lab 5 pauses a running flow until a person approves or rejects.", RED)
    card(s, Inches(6.82), Inches(4.55), Inches(4.9), Inches(1.55), "HTTP request pattern", "Labs 6A–6B let an external website POST JSON and receive JSON back.", GREEN)


def draw_trigger_map(s):
    flow(s, 2.15, ["Person\npresses Run", "Clock\nrecurs", "Form\nyields event", "Website\nPOSTs JSON", "Copilot agent\ncalls tool"], ["instant", "scheduled", "automated", "HTTP"], [BLUE, TEAL, VIOLET, GREEN, AMBER], height=1.4)
    text(s, Inches(0.72), Inches(4.25), Inches(11.85), Inches(0.48), "One trigger starts a flow. The actions may stay the same even when the trigger changes.", 19, INK, True, PP_ALIGN.CENTER)
    flow(s, 5.05, ["Trigger", "Validate data", "Business actions", "Return or notify"], ["starts", "then", "finish"], [BLUE, BLUE, TEAL, AMBER], x0=2.0, total=9.2, height=1.0)


def draw_approval(s):
    flow(s, 2.0, ["Request submitted", "Start and wait\nfor approval", "Human approves\nor rejects", "Condition reads\nOutcome"], ["pause", "decision", "resume"], [BLUE, RED, RED, VIOLET], height=1.35)
    card(s, Inches(1.15), Inches(4.25), Inches(5.25), Inches(1.45), "If yes — Approve", "Send the approved message and continue the business process.", GREEN)
    card(s, Inches(6.85), Inches(4.25), Inches(5.25), Inches(1.45), "If no — Reject", "Send the rejected message and stop or return for changes.", RED)
    text(s, Inches(0.72), Inches(6.08), Inches(11.85), Inches(0.48), "The automation is still deterministic; the person supplies the decision.", 16, INK, True, PP_ALIGN.CENTER)


def draw_day1_labs(s):
    flow(s, 2.05, ["Labs 1–2\nInstant", "Lab 3\nScheduled", "Lab 4\nAutomated", "Lab 5\nApproval", "Labs 6A–B\nHTTP"], ["add time", "add event", "add person", "add external caller"], [BLUE, TEAL, VIOLET, RED, GREEN], height=1.35)
    text(s, Inches(0.72), Inches(4.05), Inches(11.85), Inches(0.42), "DAY 1 BUILDS THE AUTOMATION ENGINE", 14, BLUE, True, PP_ALIGN.CENTER)
    card(s, Inches(1.0), Inches(4.65), Inches(3.45), Inches(1.45), "Trigger", "What starts the process?", BLUE)
    card(s, Inches(4.95), Inches(4.65), Inches(3.45), Inches(1.45), "Actions", "What work should happen?", TEAL)
    card(s, Inches(8.9), Inches(4.65), Inches(3.45), Inches(1.45), "Outputs", "What result must be stored or returned?", AMBER)


def draw_lab_instant(s):
    flow(s, 2.15, ["Learner prompt", "Copilot draft", "Manual trigger", "Send email", "Inbox"], ["generate", "review", "dynamic data", "deliver"], [VIOLET, VIOLET, BLUE, TEAL, AMBER], height=1.35)
    card(s, Inches(1.0), Inches(4.35), Inches(5.25), Inches(1.45), "Concept", "An instant flow runs only when a person deliberately starts it.", BLUE)
    card(s, Inches(6.75), Inches(4.35), Inches(5.25), Inches(1.45), "Workplace use", "A service officer acknowledges a verified phone or counter enquiry within 15 minutes.", TEAL)


def draw_lab_excel(s):
    flow(s, 2.15, ["Press Run", "Enter enquiry fields", "Map table columns", "Add Excel row", "Verify run history"], ["inputs", "dynamic content", "write", "check"], [BLUE, BLUE, TEAL, GREEN, AMBER], height=1.35)
    card(s, Inches(1.0), Inches(4.35), Inches(5.25), Inches(1.45), "Concept", "Actions consume outputs from the trigger as structured dynamic content.", BLUE)
    card(s, Inches(6.75), Inches(4.35), Inches(5.25), Inches(1.45), "Workplace use", "The shared enquiry register gives supervisors one visible queue and a basic audit trail.", GREEN)


def draw_lab_scheduled(s):
    flow(s, 2.15, ["Recurrence", "Weekdays", "Singapore time", "Send reminder", "Team inbox"], ["configure", "select", "fire", "deliver"], [TEAL, TEAL, BLUE, VIOLET, AMBER], height=1.35)
    card(s, Inches(1.0), Inches(4.35), Inches(5.25), Inches(1.45), "Workplace use", "At 9 AM on weekdays, the team lead reminds officers to clear records still marked New.", TEAL)
    card(s, Inches(6.75), Inches(4.35), Inches(5.25), Inches(1.45), "Critical setting", "Always set the business time zone; recurrence defaults can otherwise use UTC.", RED)


def draw_lab_automated(s):
    flow(s, 2.05, ["Customer submits\nMicrosoft Form", "New response\ntrigger", "Get response\ndetails", "Send email", "Add Excel row"], ["event", "Response Id", "fan out", "log"], [VIOLET, VIOLET, BLUE, TEAL, GREEN], height=1.5)
    card(s, Inches(1.0), Inches(4.5), Inches(5.25), Inches(1.45), "Workplace use", "A customer form creates the service email and audit row without staff re-keying.", VIOLET)
    card(s, Inches(6.75), Inches(4.5), Inches(5.25), Inches(1.45), "Data pattern", "Get response details turns the response ID into fields that email and Excel can use.", BLUE)


def draw_lab_approval_http(s):
    card(s, Inches(0.72), Inches(1.9), Inches(3.65), Inches(4.7), "Lab 5 — Laptop approval", "SGD 1,850 request\n↓\nManager decides\n↓\nApprove / Reject\n↓\nNotify requester", RED, 5, 15)
    card(s, Inches(4.72), Inches(1.9), Inches(3.65), Inches(4.7), "Lab 6A — Website enquiry", "Business-banking form\n↓\nPOST JSON\n↓\nNotify operations\n↓\nConfirm on page", GREEN, "6A", 15)
    card(s, Inches(8.72), Inches(1.9), Inches(3.65), Inches(4.7), "Lab 6B — Service widget", "Customer question\n↓\nApproved routes\n↓\nSafe fallback\n↓\nReturn bot reply", BLUE, "6B", 15)


def draw_day1_recap(s):
    rows = [
        ("Instant", "Person presses Run", "Labs 1–2", "Learn actions and data mapping"),
        ("Scheduled", "Recurrence", "Lab 3", "Run by the clock"),
        ("Automated", "Forms event", "Lab 4", "React to a business event"),
        ("Approval", "Flow pauses for a person", "Lab 5", "Add human judgement"),
        ("HTTP", "External JSON request", "Labs 6A–B", "Connect a website or system"),
    ]
    for i, (kind, trigger, lab, purpose) in enumerate(rows):
        y = Inches(1.82 + i * 0.94)
        box(s, Inches(0.72), y, Inches(11.85), Inches(0.75), LIGHT, LINE)
        text(s, Inches(0.92), y, Inches(1.55), Inches(0.75), kind, 15, PALETTE[i], True)
        text(s, Inches(2.55), y, Inches(3.0), Inches(0.75), trigger, 13, INK)
        text(s, Inches(5.75), y, Inches(1.35), Inches(0.75), lab, 13, BLUE, True, PP_ALIGN.CENTER)
        text(s, Inches(7.25), y, Inches(5.0), Inches(0.75), purpose, 13, GREY)


def draw_day2_roadmap(s):
    flow(s, 2.05, ["Lab 7A\nAgent shell", "Lab 7B\nRAG knowledge", "Lab 8\nTeams + website", "Lab 9\nAgent flow", "Lab 10\nPrompt flow"], ["ground", "publish", "take action", "generate"], [VIOLET, BLUE, TEAL, GREEN, AMBER], height=1.5)
    text(s, Inches(0.72), Inches(4.2), Inches(11.85), Inches(0.42), "DAY 2 ADDS THE CONVERSATIONAL ORCHESTRATOR", 14, VIOLET, True, PP_ALIGN.CENTER)
    card(s, Inches(1.0), Inches(4.75), Inches(3.45), Inches(1.35), "Understand", "Instructions and confirmed user intent", VIOLET)
    card(s, Inches(4.95), Inches(4.75), Inches(3.45), Inches(1.35), "Ground", "Retrieve approved facts with RAG", BLUE)
    card(s, Inches(8.9), Inches(4.75), Inches(3.45), Inches(1.35), "Act", "Call Power Automate as a tool", GREEN)


def draw_agent_flow_arch(s):
    flow(s, 2.0, ["User", "Copilot Studio\nagent", "Power Automate\nagent flow", "Business systems", "Respond to\nthe agent"], ["chat", "tool call", "actions", "outputs"], [TEAL, VIOLET, GREEN, BLUE, AMBER], height=1.45)
    card(s, Inches(0.85), Inches(4.35), Inches(3.5), Inches(1.45), "Agent = brain and mouth", "Understands, asks, confirms and chooses a tool.", VIOLET)
    card(s, Inches(4.9), Inches(4.35), Inches(3.5), Inches(1.45), "Flow = hands", "Logs, emails, approves, calls APIs and applies rules.", GREEN)
    card(s, Inches(8.95), Inches(4.35), Inches(3.5), Inches(1.45), "Return value = answer", "Reference, decision or draft comes back into chat.", AMBER)


def draw_n8n_mapping(s):
    mappings = [
        ("Chat Trigger", "Teams / website channel"),
        ("AI Agent node", "Copilot Studio agent"),
        ("System prompt", "Agent Instructions"),
        ("Vector / RAG tool", "Agent Knowledge"),
        ("Workflow tool", "Power Automate agent flow"),
        ("Tool result", "Respond to the agent"),
    ]
    text(s, Inches(0.88), Inches(1.82), Inches(4.9), Inches(0.35), "n8n AI-agent pattern", 17, INK, True, PP_ALIGN.CENTER)
    text(s, Inches(7.48), Inches(1.82), Inches(4.9), Inches(0.35), "Microsoft equivalent", 17, INK, True, PP_ALIGN.CENTER)
    for i, (left, right) in enumerate(mappings):
        y = Inches(2.35 + i * 0.66)
        box(s, Inches(0.88), y, Inches(4.9), Inches(0.5), LIGHT, LINE)
        text(s, Inches(1.05), y, Inches(4.55), Inches(0.5), left, 12, GREY, i == 1)
        arrow(s, Inches(6.18), y + Inches(0.07), Inches(0.7), Inches(0.36), PALETTE[i % 5])
        box(s, Inches(7.48), y, Inches(4.9), Inches(0.5), WHITE, PALETTE[i % 5])
        text(s, Inches(7.65), y, Inches(4.55), Inches(0.5), right, 12, INK, i == 1)


def draw_day2_labs(s):
    flow(s, 2.0, ["7A\nCreate agent", "7B\nAdd RAG", "8\nDeploy channels", "9\nAgent flow", "10\nPrompt flow"], ["facts", "reach users", "deterministic action", "controlled AI"], [VIOLET, BLUE, TEAL, GREEN, AMBER], height=1.55)
    card(s, Inches(0.95), Inches(4.45), Inches(3.55), Inches(1.45), "Build or configure", "Create the agent and inspect every generated instruction.", VIOLET)
    card(s, Inches(4.9), Inches(4.45), Inches(3.55), Inches(1.45), "Test evidence", "Use positive, negative and unsupported questions.", BLUE)
    card(s, Inches(8.85), Inches(4.45), Inches(3.55), Inches(1.45), "Publish safely", "Validate in Preview before Teams or website channels.", TEAL)


def draw_lab_agent(s):
    flow(s, 2.05, ["Natural-language\nrequirement", "Generated agent", "Review\nInstructions", "Preview and test", "Refine and save"], ["create", "inspect", "exercise", "improve"], [VIOLET, VIOLET, BLUE, TEAL, AMBER], height=1.45)
    card(s, Inches(1.0), Inches(4.45), Inches(5.25), Inches(1.4), "Workplace role", "First-line IT support for password, MFA, VPN and service-desk escalation.", VIOLET)
    card(s, Inches(6.75), Inches(4.45), Inches(5.25), Inches(1.4), "What it does not have yet", "No approved internal facts and no Power Automate tool.", GREY)


def draw_lab_rag(s):
    flow(s, 2.05, ["User question", "Agent searches", "Approved IT FAQ", "Retrieve passage", "Grounded answer\n+ citation"], ["query", "knowledge", "evidence", "generate"], [TEAL, VIOLET, BLUE, BLUE, GREEN], height=1.45)
    card(s, Inches(1.0), Inches(4.45), Inches(5.25), Inches(1.4), "Evidence found", "Give the approved VPN, password or security procedure and cite the FAQ.", GREEN)
    card(s, Inches(6.75), Inches(4.45), Inches(5.25), Inches(1.4), "No authority", "Refuse access approvals and route the employee to the Service Desk.", RED)


def draw_lab_channels(s):
    card(s, Inches(0.72), Inches(1.9), Inches(5.65), Inches(4.55), "Part A — Copilot agent in Teams", "User chats in Teams\n↓\nCopilot Studio agent\n↓\nGrounded conversation\n↓\nAnswer in Teams", VIOLET, "A", 16)
    card(s, Inches(6.92), Inches(1.9), Inches(5.65), Inches(4.55), "Part B — Website calls deterministic flow", "Website enquiry form\n↓\nHTTP POST\n↓\nDeterministic Power Automate\n↓\nJSON result on page", GREEN, "B", 16)
    text(s, Inches(0.72), Inches(6.55), Inches(11.85), Inches(0.28), "Comparison: Part B is automation only—the agent is not in that submission path yet.", 13, RED, True, PP_ALIGN.CENTER)


def draw_lab_agent_flow(s):
    flow(s, 2.0, ["Teams or website\nchat", "Copilot agent\nconfirms inputs", "Agent flow tool", "Deterministic\nconditions", "Respond to\nagent", "Answer in chat"], ["chat", "call", "evaluate", "outputs", "present"], [TEAL, VIOLET, GREEN, BLUE, AMBER, TEAL], height=1.4)
    card(s, Inches(1.0), Inches(4.35), Inches(5.25), Inches(1.45), "Same as n8n AI Agent + workflow tool", "The agent decides when to call the tool; the flow executes predictable business rules.", VIOLET)
    card(s, Inches(6.75), Inches(4.35), Inches(5.25), Inches(1.45), "Required endpoints", "When an agent calls the flow → actions → Respond to the agent.", GREEN)


def draw_lab_prompt_flow(s):
    flow(s, 1.95, ["Confirmed enquiry", "Agent calls\nprompt flow", "AI Builder\nRun a prompt", "Parse + validate\nstructured JSON", "Guardrail decision", "Respond to\nagent"], ["tool", "generate", "schema", "check", "safe output"], [VIOLET, GREEN, AMBER, BLUE, RED, TEAL], height=1.45)
    card(s, Inches(1.0), Inches(4.35), Inches(5.25), Inches(1.5), "Use AI where language varies", "Draft, classify or summarise—but keep approvals, identity checks and limits deterministic.", AMBER)
    card(s, Inches(6.75), Inches(4.35), Inches(5.25), Inches(1.5), "Validate before returning", "Parse the schema, check safety rules and fall back when the output is missing or unsafe.", RED)


def draw_compare(s):
    columns = [
        ("Lab 8", "Ordinary flow", "Website HTTP request", "No agent in path", "JSON to webpage", GREEN),
        ("Lab 9", "Agent flow", "Agent tool call", "Deterministic rules", "Structured result to chat", VIOLET),
        ("Lab 10", "Prompt flow", "Agent tool call", "AI prompt + guardrails", "Controlled draft to chat", AMBER),
    ]
    for i, (lab, kind, trigger, processing, result, color) in enumerate(columns):
        x = Inches(0.72 + i * 4.0)
        card(s, x, Inches(1.9), Inches(3.7), Inches(4.65), f"{lab} — {kind}", f"TRIGGER\n{trigger}\n\nPROCESSING\n{processing}\n\nRETURN\n{result}", color, i + 8, 13)


def draw_day2_recap(s):
    flow(s, 1.95, ["Instructions", "Knowledge (RAG)", "Channels", "Tools", "Prompt + guardrails"], ["ground", "publish", "act", "generate"], [VIOLET, BLUE, TEAL, GREEN, AMBER], height=1.35)
    text(s, Inches(0.72), Inches(3.75), Inches(11.85), Inches(0.48), "THE FINAL ARCHITECTURE", 14, VIOLET, True, PP_ALIGN.CENTER)
    flow(s, 4.45, ["User", "Copilot agent", "Power Automate tool", "Business result", "Answer in chat"], ["asks", "calls", "returns", "presents"], [TEAL, VIOLET, GREEN, BLUE, AMBER], height=1.3)


def draw_big_picture(s):
    card(s, Inches(0.72), Inches(1.85), Inches(3.55), Inches(4.7), "1 — Automate", "Instant\nScheduled\nAutomated\nApproval\nHTTP", BLUE, 1, 17)
    arrow(s, Inches(4.47), Inches(3.75), Inches(0.58), Inches(0.5), BLUE)
    card(s, Inches(5.15), Inches(1.85), Inches(3.55), Inches(4.7), "2 — Add the agent", "Instructions\nRAG knowledge\nTeams + website\nConfirmed inputs", VIOLET, 2, 17)
    arrow(s, Inches(8.9), Inches(3.75), Inches(0.58), Inches(0.5), VIOLET)
    card(s, Inches(9.58), Inches(1.85), Inches(3.0), Inches(4.7), "3 — Let it act", "Agent flow\nPrompt flow\nGuardrails\nReturned result", GREEN, 3, 17)


rebuild(7, "What You’ll Learn", "Course outcomes", draw_course_outcomes)
rebuild(8, "Your 2-Day Journey", "Simple to complex", draw_journey)
rebuild(22, "Triggers — What Starts a Workflow", "One start event", draw_trigger_map)
rebuild(30, "Types of Flow and Process Patterns", "Power Automate", draw_trigger_spectrum)
rebuild(39, "Human in the Loop — Approval Pattern", "Pause · decide · resume", draw_approval)
rebuild(41, "Hands-On Labs — Power Automate", "Day 1 progression", draw_day1_labs)
rebuild(42, "Lab 1 — Instant Email Flow", "Prompt → review → test", draw_lab_instant)
rebuild(43, "Lab 2 — Instant Excel Logging Flow", "Inputs → dynamic content → table", draw_lab_excel)
rebuild(44, "Lab 3 — Scheduled Flow", "Recurrence and time zone", draw_lab_scheduled)
rebuild(45, "Lab 4 — Automated Form Flow", "Event-driven automation", draw_lab_automated)
rebuild(46, "Labs 5–6 — Human and External Triggers", "Approval and HTTP", draw_lab_approval_http)
rebuild(47, "Day 1 Recap — Choose the Right Trigger", "Automation foundations", draw_day1_recap)
rebuild(49, "Day 2 Roadmap", "Agent capabilities build progressively", draw_day2_roadmap)
rebuild(60, "Connecting Agents to Power Automate", "Agent calls a flow as a tool", draw_agent_flow_arch)
rebuild(64, "n8n AI Agent Pattern in Microsoft", "Concept mapping", draw_n8n_mapping)
rebuild(65, "Hands-On Labs — Copilot Studio", "Day 2 progression", draw_day2_labs)
rebuild(66, "Lab 7A — Create the IT Support Agent", "Instructions before knowledge", draw_lab_agent)
rebuild(67, "Lab 7B — Ground the Agent with RAG", "Retrieve before generating", draw_lab_rag)
rebuild(68, "Lab 8 — Deploy to Teams and Website", "Two channels · two paths", draw_lab_channels)
rebuild(69, "Lab 9 — Agent Calls a Deterministic Flow", "Copilot + Power Automate", draw_lab_agent_flow)
rebuild(70, "Lab 10 — Agent Calls a Prompt Flow", "AI Builder with guardrails", draw_lab_prompt_flow)
rebuild(71, "Compare Labs 8, 9 and 10", "Automation → agent flow → prompt flow", draw_compare)
rebuild(73, "Day 2 Recap — The Agentic Architecture", "Understand · ground · act", draw_day2_recap)
rebuild(79, "Course Recap — The Big Picture", "From flow maker to agent builder", draw_big_picture)

prs.save(DECK)
print(f"Updated {DECK} ({len(prs.slides)} slides)")
